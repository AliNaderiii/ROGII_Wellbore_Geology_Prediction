"""Section 2 — Task presentation audit.

Extracts slide titles, body text, tables, speaker notes and embedded images
from the official task .pptx into reports/task_presentation_summary.md.

The extractor is purely mechanical: it transcribes what the deck says. It
does not add interpretation. A short "open questions" section lists things
the deck does NOT answer, so nothing is silently invented downstream.
"""
from __future__ import annotations

import sys
import zipfile
from pathlib import Path

try:
    from scripts._bootstrap import bootstrap
except ImportError:  # executed as a loose file, not as a package
    sys.path.insert(0, str(Path.cwd()))
    from scripts._bootstrap import bootstrap
bootstrap()

from src.paths import REPORTS_DIR, TASK_PPTX, ensure_reports_dir



def extract_with_pptx(path: Path) -> list[str]:
    from pptx import Presentation

    prs = Presentation(str(path))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()
        lines.append(f"\n## Slide {i}" + (f" — {title}" if title else ""))
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and text != title:
                    for para in text.splitlines():
                        if para.strip():
                            lines.append(f"- {para.strip()}")
            if getattr(shape, "has_table", False) and shape.has_table:
                tbl = shape.table
                lines.append("")
                for r, row in enumerate(tbl.rows):
                    cells = [c.text.strip().replace("|", "/") for c in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
                    if r == 0:
                        lines.append("|" + "---|" * len(cells))
                lines.append("")
            if shape.shape_type == 13:  # PICTURE
                try:
                    img = shape.image
                    (REPORTS_DIR / 'task_presentation_images').mkdir(parents=True, exist_ok=True)
                    name = f"slide{i:02d}_{shape.shape_id}.{img.ext}"
                    (REPORTS_DIR / 'task_presentation_images' / name).write_bytes(img.blob)
                    lines.append(f"- ![embedded image](task_presentation_images/{name})")
                except Exception as exc:  # pragma: no cover
                    lines.append(f"- (image extraction failed: {exc})")
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"\n> **Speaker notes:** {notes}")
        except Exception:
            pass
    return lines


def extract_with_zip(path: Path) -> list[str]:
    """Dependency-free fallback: pull raw <a:t> runs out of the OOXML."""
    import re
    import xml.etree.ElementTree as ET

    lines: list[str] = []
    with zipfile.ZipFile(path) as zf:
        slides = sorted(
            (n for n in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", n)),
            key=lambda n: int(re.findall(r"\d+", n)[-1]),
        )
        for i, name in enumerate(slides, start=1):
            root = ET.fromstring(zf.read(name))
            texts = [
                (e.text or "").strip()
                for e in root.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}t")
            ]
            texts = [t for t in texts if t]
            lines.append(f"\n## Slide {i}" + (f" — {texts[0]}" if texts else ""))
            for t in texts[1:]:
                lines.append(f"- {t}")
    return lines


def main() -> None:
    ensure_reports_dir()
    if not TASK_PPTX.exists():
        raise SystemExit(f"Presentation not found: {TASK_PPTX}")

    try:
        body = extract_with_pptx(TASK_PPTX)
        engine = "python-pptx (text, tables, notes, embedded images)"
    except ImportError:
        body = extract_with_zip(TASK_PPTX)
        engine = "stdlib zipfile + ElementTree fallback (text only)"

    header = [
        "# Task Presentation Summary",
        "",
        f"Source: `{TASK_PPTX}`",
        f"Extraction engine: {engine}",
        "",
        "This file is a **verbatim transcription** of the official deck. Nothing",
        "below is inferred; any modelling opinion lives in the other reports.",
        "",
        "---",
    ]
    footer = [
        "",
        "---",
        "",
        "## Points the deck must be checked against before modelling",
        "",
        "- [ ] Exact target definition (TVT units, sign convention, datum).",
        "- [ ] Evaluation metric and whether it is per-row or per-well averaged.",
        "- [ ] Whether external data / pretrained artifacts are permitted.",
        "- [ ] Whether the submission must cover only hidden-suffix rows.",
        "- [ ] Any runtime / notebook-only constraint.",
        "",
        "Tick each box against the transcription above and the official rules page",
        "before any of the external artifacts are promoted from",
        "'NEEDS FURTHER REVIEW' to 'USE'.",
    ]
    OUT = REPORTS_DIR / 'task_presentation_summary.md'
    OUT.write_text("\n".join(header + body + footer), encoding="utf-8")
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
