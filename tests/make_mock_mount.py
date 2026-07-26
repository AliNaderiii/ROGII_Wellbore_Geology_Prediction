"""Build a small synthetic /kaggle-style mount so the audit + pipeline code can
be exercised without the real competition data.

    python tests/make_mock_mount.py /tmp/mockkaggle

This is a TEST FIXTURE ONLY. It contains no competition data.
"""
from __future__ import annotations

import sys
import zipfile
from pathlib import Path

import numpy as np
import pandas as pd

FORMATIONS = ["ANCC", "ASTNU", "ASTNL", "EGFDU", "EGFDL", "BUDA"]
RNG = np.random.default_rng(0)


def make_well(well_id: str, hw_dir: Path, *, with_target: bool, with_markers: bool, n=400):
    md = np.arange(n, dtype=float) * 0.5 + 8000.0
    tvt = np.cumsum(RNG.normal(0, 0.05, n)) + 20.0
    gr = 60 + 25 * np.sin(md / 40) + RNG.normal(0, 5, n)
    gr[50:90] = np.nan  # contiguous outage
    prefix = int(n * 0.25)
    tvt_input = tvt.copy()
    tvt_input[prefix:] = np.nan

    df = pd.DataFrame({
        "MD": md,
        "X": 1000 + md * 0.3,
        "Y": 2000 + md * 0.1,
        "Z": -md * 0.9,
        "GR": gr,
        "TVT_input": tvt_input,
    })
    if with_markers:
        for i, f in enumerate(FORMATIONS):
            df[f] = 10.0 + 8 * i + RNG.normal(0, 0.1, n)
    if with_target:
        df["TVT"] = tvt
    df.to_csv(hw_dir / f"{well_id}__horizontal_well.csv", index=False)

    tw_tvt = np.linspace(0, 60, 300)
    tw = pd.DataFrame({
        "TVT": tw_tvt,
        "GR": 50 + 30 * np.sin(tw_tvt / 6),
        "Geology": [FORMATIONS[min(int(t // 10), 5)] for t in tw_tvt],
    })
    tw.to_csv(hw_dir / f"{well_id}__typewell.csv", index=False)
    (hw_dir / f"{well_id}__log.png").write_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * 40)
    return prefix, n


def main(root: Path) -> None:
    comp = root / "input" / "competitions" / "rogii-wellbore-geology-prediction"
    train, test = comp / "train", comp / "test"
    train.mkdir(parents=True, exist_ok=True)
    test.mkdir(parents=True, exist_ok=True)

    for i in range(5):
        make_well(f"TRW{i:03d}", train, with_target=True, with_markers=True)

    rows = []
    for i in range(3):
        p, n = make_well(f"TSW{i:03d}", test, with_target=False, with_markers=False)
        rows += [{"id": f"TSW{i:03d}_{r}", "tvt": 0.0} for r in range(p, n)]
    pd.DataFrame(rows).to_csv(comp / "sample_submission.csv", index=False)

    # minimal pptx
    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = "Wellbore Geology Prediction"
        s.placeholders[1].text = "Predict TVT along the hidden suffix.\nMetric: RMSE."
        prs.save(comp / "AI_wellbore_geology_prediction_task_en.pptx")
    except ImportError:
        (comp / "AI_wellbore_geology_prediction_task_en.pptx").write_bytes(b"")

    ds = root / "input" / "datasets"
    kb = ds / "phongnguyn23021656" / "koolbox-offline"
    kb.mkdir(parents=True, exist_ok=True)
    whl = kb / "koolbox-0.1.2-py3-none-any.whl"
    with zipfile.ZipFile(whl, "w") as zf:
        zf.writestr("koolbox-0.1.2.dist-info/METADATA",
                    "Name: koolbox\nVersion: 0.1.2\nRequires-Python: >=3.8\nLicense: MIT\n")

    cm = ds / "fleongg" / "rogii-claude-models-pub"
    cm.mkdir(parents=True, exist_ok=True)
    import pickle
    (cm / "model_fold0.pkl").write_bytes(pickle.dumps({"a": np.zeros(3)}))
    (cm / "README.md").write_text("Trained on train wells. LightGBM.\n")

    ar = ds / "ravaghi" / "wellbore-geology-prediction-artifacts"
    ar.mkdir(parents=True, exist_ok=True)
    pd.DataFrame({"well_id": ["TRW000"], "tvt": [1.0]}).to_csv(ar / "oof_predictions.csv", index=False)
    pd.DataFrame({"id": ["TSW000_100"], "tvt": [1.0]}).to_csv(ar / "submission_blend.csv", index=False)

    print(f"mock mount at {root}")


if __name__ == "__main__":
    main(Path(sys.argv[1] if len(sys.argv) > 1 else "/tmp/mockkaggle"))
